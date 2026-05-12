#!/usr/bin/env python3
"""Assemble final full-slide images into a PPTX deck.

Uses python-pptx when available. Falls back to a minimal standard-library
OOXML writer so final deck assembly does not depend on optional packages.
"""

from __future__ import annotations

import argparse
import re
import zipfile
from datetime import datetime, timezone
from html import escape
from pathlib import Path


def optional_pptx():
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches  # type: ignore
    except Exception as exc:  # pragma: no cover - depends on local environment
        return None, None, exc
    return Presentation, Inches


def slide_number(path: Path) -> int:
    match = re.search(r"slide-(\d+)", path.stem)
    return int(match.group(1)) if match else 0


def content_type_for(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".png":
        return "image/png"
    if suffix in {".jpg", ".jpeg"}:
        return "image/jpeg"
    raise ValueError(f"Unsupported image type: {path}")


def image_extension(path: Path) -> str:
    suffix = path.suffix.lower().lstrip(".")
    return "jpeg" if suffix == "jpg" else suffix


def write_pptx_with_python_pptx(project_dir: Path, images: list[Path], output: Path) -> None:
    Presentation, Inches, _ = optional_pptx()
    if Presentation is None or Inches is None:
        raise RuntimeError("python-pptx unavailable")

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for image in images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(image, 0, 0, width=prs.slide_width, height=prs.slide_height)

    prs.save(output)


def write_minimal_pptx(project_dir: Path, images: list[Path], output: Path) -> None:
    width = 12192000
    height = 6858000
    now = datetime.now(timezone.utc).replace(microsecond=0).isoformat().replace("+00:00", "Z")

    slide_overrides = "\n".join(
        f'<Override PartName="/ppt/slides/slide{i}.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
        for i in range(1, len(images) + 1)
    )
    image_defaults = "\n".join(
        sorted(
            {
                f'<Default Extension="{image_extension(image)}" ContentType="{content_type_for(image)}"/>'
                for image in images
            }
        )
    )
    presentation_slide_ids = "\n".join(
        f'<p:sldId id="{255 + i}" r:id="rId{i + 1}"/>'
        for i in range(1, len(images) + 1)
    )
    presentation_rels = "\n".join(
        [('<Relationship Id="rId1" '
          'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" '
          'Target="slideMasters/slideMaster1.xml"/>')]
        + [
            f'<Relationship Id="rId{i + 1}" '
            'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" '
            f'Target="slides/slide{i}.xml"/>'
            for i in range(1, len(images) + 1)
        ]
    )

    content_types = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  {image_defaults}
  <Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>
  <Override PartName="/docProps/app.xml" ContentType="application/vnd.openxmlformats-officedocument.extended-properties+xml"/>
  <Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
  <Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>
  <Override PartName="/ppt/slideLayouts/slideLayout1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>
  <Override PartName="/ppt/theme/theme1.xml" ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>
  {slide_overrides}
</Types>'''

    root_rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/>
  <Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties" Target="docProps/app.xml"/>
</Relationships>'''

    core = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/" xmlns:dcterms="http://purl.org/dc/terms/" xmlns:dcmitype="http://purl.org/dc/dcmitype/" xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">
  <dc:title>paper-to-ppt generated deck</dc:title>
  <dc:creator>paper-to-ppt</dc:creator>
  <cp:lastModifiedBy>paper-to-ppt</cp:lastModifiedBy>
  <dcterms:created xsi:type="dcterms:W3CDTF">{now}</dcterms:created>
  <dcterms:modified xsi:type="dcterms:W3CDTF">{now}</dcterms:modified>
</cp:coreProperties>'''

    app = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties" xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">
  <Application>paper-to-ppt</Application>
  <PresentationFormat>On-screen Show (16:9)</PresentationFormat>
  <Slides>{len(images)}</Slides>
</Properties>'''

    presentation = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:sldMasterIdLst><p:sldMasterId id="2147483648" r:id="rId1"/></p:sldMasterIdLst>
  <p:sldIdLst>{presentation_slide_ids}</p:sldIdLst>
  <p:sldSz cx="{width}" cy="{height}" type="wide"/>
  <p:notesSz cx="6858000" cy="9144000"/>
  <p:defaultTextStyle/>
</p:presentation>'''

    presentation_rels_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  {presentation_rels}
</Relationships>'''

    theme = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="paper-to-ppt">
  <a:themeElements>
    <a:clrScheme name="paper-to-ppt">
      <a:dk1><a:srgbClr val="111111"/></a:dk1><a:lt1><a:srgbClr val="FFFFFF"/></a:lt1>
      <a:dk2><a:srgbClr val="222222"/></a:dk2><a:lt2><a:srgbClr val="F7F7F7"/></a:lt2>
      <a:accent1><a:srgbClr val="3366CC"/></a:accent1><a:accent2><a:srgbClr val="DC3912"/></a:accent2>
      <a:accent3><a:srgbClr val="FF9900"/></a:accent3><a:accent4><a:srgbClr val="109618"/></a:accent4>
      <a:accent5><a:srgbClr val="990099"/></a:accent5><a:accent6><a:srgbClr val="0099C6"/></a:accent6>
      <a:hlink><a:srgbClr val="0000FF"/></a:hlink><a:folHlink><a:srgbClr val="800080"/></a:folHlink>
    </a:clrScheme>
    <a:fontScheme name="paper-to-ppt"><a:majorFont><a:latin typeface="Aptos Display"/></a:majorFont><a:minorFont><a:latin typeface="Aptos"/></a:minorFont></a:fontScheme>
    <a:fmtScheme name="paper-to-ppt"><a:fillStyleLst/><a:lnStyleLst/><a:effectStyleLst/><a:bgFillStyleLst/></a:fmtScheme>
  </a:themeElements>
</a:theme>'''

    slide_master = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld><p:bg><p:bgPr><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></p:bgPr></p:bg><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr></p:spTree></p:cSld>
  <p:sldLayoutIdLst><p:sldLayoutId id="2147483649" r:id="rId1"/></p:sldLayoutIdLst>
  <p:txStyles><p:titleStyle/><p:bodyStyle/><p:otherStyle/></p:txStyles>
</p:sldMaster>'''

    slide_master_rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/>
</Relationships>'''

    slide_layout = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldLayout xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" type="blank" preserve="1">
  <p:cSld name="Blank"><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr></p:spTree></p:cSld>
</p:sldLayout>'''

    slide_layout_rels = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="../slideMasters/slideMaster1.xml"/>
</Relationships>'''

    def slide_xml(index: int, image_name: str) -> str:
        image_name = escape(image_name)
        return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
      <p:pic>
        <p:nvPicPr><p:cNvPr id="{index + 1}" name="{image_name}"/><p:cNvPicPr/><p:nvPr/></p:nvPicPr>
        <p:blipFill><a:blip r:embed="rId1"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>
        <p:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="{width}" cy="{height}"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>
      </p:pic>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sld>'''

    def slide_rels(index: int, image_ext: str) -> str:
        return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/image{index}.{image_ext}"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
</Relationships>'''

    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", content_types)
        zf.writestr("_rels/.rels", root_rels)
        zf.writestr("docProps/core.xml", core)
        zf.writestr("docProps/app.xml", app)
        zf.writestr("ppt/presentation.xml", presentation)
        zf.writestr("ppt/_rels/presentation.xml.rels", presentation_rels_xml)
        zf.writestr("ppt/theme/theme1.xml", theme)
        zf.writestr("ppt/slideMasters/slideMaster1.xml", slide_master)
        zf.writestr("ppt/slideMasters/_rels/slideMaster1.xml.rels", slide_master_rels)
        zf.writestr("ppt/slideLayouts/slideLayout1.xml", slide_layout)
        zf.writestr("ppt/slideLayouts/_rels/slideLayout1.xml.rels", slide_layout_rels)
        for index, image in enumerate(images, start=1):
            ext = image_extension(image)
            zf.write(image, f"ppt/media/image{index}.{ext}")
            zf.writestr(f"ppt/slides/slide{index}.xml", slide_xml(index, image.name))
            zf.writestr(f"ppt/slides/_rels/slide{index}.xml.rels", slide_rels(index, ext))


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("project_dir", help="Project directory.")
    parser.add_argument("--renders-dir", default="renders/final", help="Directory of slide images relative to project.")
    parser.add_argument("--output", default="deck/final_deck.pptx", help="Output PPTX relative to project.")
    parser.add_argument("--wide", action="store_true", default=True, help="Use 16:9 widescreen dimensions.")
    args = parser.parse_args()

    project_dir = Path(args.project_dir).expanduser().resolve()
    renders_dir = project_dir / args.renders_dir
    output = project_dir / args.output
    output.parent.mkdir(parents=True, exist_ok=True)

    images = sorted(
        [p for p in renders_dir.iterdir() if p.suffix.lower() in {".png", ".jpg", ".jpeg"}],
        key=slide_number,
    )
    if not images:
        raise SystemExit(f"No slide images found in {renders_dir}")

    try:
        write_pptx_with_python_pptx(project_dir, images, output)
        writer = "python-pptx"
    except RuntimeError:
        write_minimal_pptx(project_dir, images, output)
        writer = "stdlib-ooxml"

    print(f"Wrote {output} with {len(images)} slides using {writer}.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
