#!/usr/bin/env python3
"""Inspect a PPTX style reference and emit layout, color, and typography statistics."""

from __future__ import annotations

import argparse
import json
from collections import Counter
from pathlib import Path


def require_pptx():
    try:
        from pptx import Presentation  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    except Exception as exc:  # pragma: no cover - depends on local environment
        raise SystemExit(
            "python-pptx is required for PPTX inspection. Install it with: pip install python-pptx"
        ) from exc
    return Presentation, MSO_SHAPE_TYPE


def rgb_to_hex(rgb) -> str | None:
    if rgb is None:
        return None
    try:
        return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
    except Exception:
        return str(rgb)


def shape_fill_hex(shape) -> str | None:
    try:
        fill = shape.fill
        if not fill or fill.type is None:
            return None
        color = fill.fore_color
        return rgb_to_hex(color.rgb)
    except Exception:
        return None


def collect_text(shape) -> tuple[list[str], list[str], list[float]]:
    texts: list[str] = []
    fonts: list[str] = []
    sizes: list[float] = []
    if not getattr(shape, "has_text_frame", False):
        return texts, fonts, sizes
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = run.text.strip()
            if text:
                texts.append(text)
            font = run.font
            if font.name:
                fonts.append(font.name)
            if font.size:
                sizes.append(round(font.size.pt, 1))
    return texts, fonts, sizes


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("project_dir", help="Project directory created by create_project.py.")
    args = parser.parse_args()

    project_dir = Path(args.project_dir).expanduser().resolve()
    sources_dir = project_dir / "sources"
    pptx_files = sorted(sources_dir.glob("*.pptx"))
    if not pptx_files:
        report = {
            "status": "unsupported_reference_format",
            "message": "No PPTX found in sources. Infer style manually from screenshots or other provided files.",
            "source_files": [p.name for p in sorted(sources_dir.iterdir())],
        }
        out = project_dir / "analysis" / "style_inspection.json"
        out.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
        print(json.dumps(report, ensure_ascii=False, indent=2))
        return 0

    Presentation, MSO_SHAPE_TYPE = require_pptx()
    pptx_path = pptx_files[0]
    prs = Presentation(pptx_path)

    fonts: Counter[str] = Counter()
    font_sizes: Counter[str] = Counter()
    fills: Counter[str] = Counter()
    text_lengths = []
    slide_records = []

    for idx, slide in enumerate(prs.slides, start=1):
        shape_records = []
        picture_count = 0
        for shape in slide.shapes:
            fill_hex = shape_fill_hex(shape)
            if fill_hex:
                fills[fill_hex] += 1
            texts, shape_fonts, sizes = collect_text(shape)
            for font in shape_fonts:
                fonts[font] += 1
            for size in sizes:
                font_sizes[str(size)] += 1
            if texts:
                text_lengths.append(sum(len(t) for t in texts))
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_count += 1
            shape_records.append(
                {
                    "shape_type": str(shape.shape_type),
                    "left": int(shape.left),
                    "top": int(shape.top),
                    "width": int(shape.width),
                    "height": int(shape.height),
                    "has_text": bool(texts),
                    "text_preview": " ".join(texts)[:160],
                    "fill": fill_hex,
                }
            )
        slide_records.append(
            {
                "slide": idx,
                "shape_count": len(slide.shapes),
                "picture_count": picture_count,
                "text_shape_count": sum(1 for s in shape_records if s["has_text"]),
                "shapes": shape_records,
            }
        )

    report = {
        "status": "ok",
        "source_pptx": str(pptx_path.relative_to(project_dir)),
        "slide_count": len(prs.slides),
        "slide_width_emu": int(prs.slide_width),
        "slide_height_emu": int(prs.slide_height),
        "aspect_ratio": round(int(prs.slide_width) / int(prs.slide_height), 4),
        "top_fonts": fonts.most_common(12),
        "top_font_sizes_pt": font_sizes.most_common(12),
        "top_fill_colors": fills.most_common(16),
        "avg_text_chars_per_text_slide": round(sum(text_lengths) / max(len(text_lengths), 1), 1),
        "slides": slide_records,
    }

    out = project_dir / "analysis" / "style_inspection.json"
    out.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
